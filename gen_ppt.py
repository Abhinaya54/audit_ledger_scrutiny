from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── New color palette: deep navy + cyan + orange accent ───────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy
NAVY_MID    = RGBColor(0x16, 0x2B, 0x5A)   # slightly lighter navy
CYAN        = RGBColor(0x00, 0xC2, 0xCB)   # bright cyan
ORANGE      = RGBColor(0xFF, 0x6B, 0x35)   # orange accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF2, 0xF5, 0xF9)
MED_GRAY    = RGBColor(0x66, 0x72, 0x80)
DARK_TEXT   = RGBColor(0x1A, 0x1A, 0x2E)
GREEN_BG    = RGBColor(0xD4, 0xED, 0xDA)
GREEN_TXT   = RGBColor(0x1A, 0x7F, 0x3C)
ROW_A       = RGBColor(0xEA, 0xF2, 0xFF)
ROW_B       = RGBColor(0xF7, 0xF9, 0xFF)
CARD_BORDER = RGBColor(0xE0, 0xE7, 0xF2)
RULE_COLS   = [
    RGBColor(0x00, 0xC2, 0xCB),
    RGBColor(0xFF, 0x6B, 0x35),
    RGBColor(0x3B, 0x82, 0xF6),
    RGBColor(0xF5, 0x9E, 0x0B),
    RGBColor(0x10, 0xB9, 0x81),
    RGBColor(0x8B, 0x5C, 0xF6),
]

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ── Helpers ───────────────────────────────────────────────────────────────────

def set_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color, line_color=None, line_pt=1.0):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    return sh


def txbox(slide, text, left, top, width, height,
          bold=False, italic=False, size=18,
          color=None, align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = DARK_TEXT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None):
    """Navy top bar with white title + cyan underline strip."""
    rect(slide, 0, 0, 13.33, 1.45, NAVY)
    rect(slide, 0, 1.45, 13.33, 0.07, CYAN)
    txbox(slide, title, 0.45, 0.18, 12.0, 0.85,
          bold=True, size=28, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.45, 0.95, 12.0, 0.4,
              size=14, color=CYAN)


def dot_bullet(slide, items, left=0.55, top=1.7, width=12.0,
               size=17, spacing=0.5, color=None):
    if color is None:
        color = DARK_TEXT
    for i, item in enumerate(items):
        is_sub = item.startswith("    ")
        dot = "  –  " if is_sub else "  ●  "
        dot_color = MED_GRAY if is_sub else CYAN
        # dot
        txbox(slide, dot.strip(), left, top + i * spacing, 0.35, 0.45,
              size=size, color=dot_color, bold=(not is_sub))
        # text
        txbox(slide, item.strip(), left + 0.32, top + i * spacing,
              width - 0.32, 0.45, size=size, color=color)


def two_col(slide, t1, items1, t2, items2, top=1.65):
    # Left column header
    rect(slide, 0.4, top, 5.9, 0.42, NAVY_MID)
    txbox(slide, t1, 0.55, top + 0.04, 5.6, 0.38,
          bold=True, size=17, color=CYAN)
    for i, item in enumerate(items1):
        txbox(slide, "●  " + item, 0.55, top + 0.55 + i * 0.46, 5.7, 0.44,
              size=16, color=DARK_TEXT)
    # Right column header
    rect(slide, 6.95, top, 5.9, 0.42, NAVY_MID)
    txbox(slide, t2, 7.1, top + 0.04, 5.6, 0.38,
          bold=True, size=17, color=CYAN)
    for i, item in enumerate(items2):
        txbox(slide, "●  " + item, 7.1, top + 0.55 + i * 0.46, 5.7, 0.44,
              size=16, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, NAVY)

# Background accent geometry
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 5.8, 13.33, 0.08, CYAN)
rect(sl, 0, 5.88, 13.33, 1.62, NAVY_MID)
# Cyan left bar
rect(sl, 0, 0, 0.12, 7.5, CYAN)
# Orange dot accent
rect(sl, 0.55, 1.7, 0.12, 0.12, ORANGE)
rect(sl, 0.55, 1.95, 0.12, 0.12, CYAN)

txbox(sl, "Internship Review Presentation",
      0.8, 2.3, 11.5, 1.3, bold=True, size=44, color=WHITE)
txbox(sl, "RCTS, IIIT-H",
      0.8, 3.65, 6.0, 0.65, size=26, color=CYAN)
txbox(sl, "Presented by  Abhinaya",
      0.8, 6.05, 7.0, 0.55, size=18, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — About the Project
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "About the Project", "Ledger Scrutiny Assistant — Audit Intelligence Suite")
dot_bullet(sl, [
    "Automates financial anomaly detection in General Ledger (GL) data",
    "Combines a synthetic data generator, rule-based engine, and ML model",
    "Helps auditors identify suspicious transactions quickly and accurately",
    "Built as a full-stack system with two core components:",
    "    Synthetic GL Generator  →  generates realistic, privacy-safe test data",
    "    Scrutiny Engine  →  flags anomalies using 6 audit rules + Isolation Forest ML",
], top=1.68, spacing=0.54)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Tasks Given vs Completed
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Tasks Assigned vs. Completed")

TASK_ROWS = [
    ("1", "Synthetic GL Data Generator (CLI + Digital Twin mode)", "Done"),
    ("2", "Rule-Based Scrutiny Engine — 6 rules (R1-R6)",          "Done"),
    ("3", "ML Anomaly Detection — Isolation Forest",               "Done"),
    ("4", "Master Pipeline (Ingest -> Rules -> ML -> Export)",     "Done"),
    ("5", "Multi-Agent System (Validator, Detector, Reporter)",    "Done"),
    ("6", "FastAPI REST Backend (Generator + Scrutiny endpoints)", "Done"),
    ("7", "React / TypeScript Frontend with live filters & charts","Done"),
]
COL_X = [0.4, 1.25, 10.6]
COL_W = [0.75, 9.2, 1.55]
HDR   = ["#", "Task", "Status"]

# Header row
rect(sl, 0.4, 1.62, 12.45, 0.44, NAVY)
for hdr, cx, cw in zip(HDR, COL_X, COL_W):
    txbox(sl, hdr, cx + 0.06, 1.65, cw - 0.1, 0.38,
          bold=True, size=15, color=WHITE)

for i, (num, task, status) in enumerate(TASK_ROWS):
    y = 2.09 + i * 0.44
    bg = ROW_A if i % 2 == 0 else ROW_B
    for j, (txt, cx, cw) in enumerate(zip([num, task, status], COL_X, COL_W)):
        fill = GREEN_BG if j == 2 else bg
        rect(sl, cx, y, cw, 0.42, fill)
        c = GREEN_TXT if j == 2 else DARK_TEXT
        txbox(sl, txt, cx + 0.06, y + 0.04, cw - 0.1, 0.36,
              size=14, color=c, bold=(j == 2))

# Total badge
rect(sl, 0.4, 7.12, 3.0, 0.32, NAVY)
txbox(sl, "7 / 7 Tasks Completed", 0.5, 7.14, 2.8, 0.28,
      bold=True, size=13, color=CYAN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Synthetic GL Generator
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Task 1 & 2: Synthetic GL Data Generator")
two_col(sl,
    "Core Generator",
    [
        "Generates realistic GL entries (50,000+ rows)",
        "Configurable fiscal year, accounts & voucher types",
        "Randomised narrations using Faker (en_IN locale)",
        "CLI: --period, --rows, --seed, --no-inject flags",
    ],
    "Digital Twin Mode",
    [
        "Mirrors real GL statistics without raw data exposure",
        "Input: JSON summary (mean, std, distributions)",
        "Preserves amount, weekday & account distributions",
        "Privacy-safe testing against real-world profiles",
    ]
)
rect(sl, 0.4, 5.55, 12.45, 0.42, NAVY_MID)
txbox(sl, "Anomaly Injection (6 pattern types):", 0.55, 5.58, 3.5, 0.36,
      bold=True, size=14, color=CYAN)
txbox(sl, "Round Numbers  |  Weekend Entries  |  Period-End Spikes  |  Weak Narrations  |  Duplicates  |  Manual Journals",
      4.1, 5.58, 8.6, 0.36, size=13, color=LIGHT_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Rule-Based Scrutiny Engine
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Task 3: Rule-Based Scrutiny Engine (R1–R6)")

RULES = [
    ("R1", "Round Numbers",   "Amount is a multiple of 1,000 / 10,000"),
    ("R2", "Weekend Entries", "Transaction posted on a Sunday"),
    ("R3", "Period-End",      "Entry within critical 5-day period-end window"),
    ("R4", "Weak Narration",  "Narration too short or uses generic placeholder"),
    ("R5", "Duplicate Check", "Same date, ledger and amount appears multiple times"),
    ("R6", "Manual Journal",  "Manual JV flagged — requires source verification"),
]

for i, (code, name, desc) in enumerate(RULES):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.47
    y = 1.62 + row * 1.88
    # Card
    rect(sl, x, y, 6.05, 1.75, LIGHT_GRAY, CARD_BORDER, 0.8)
    # Top accent strip
    rect(sl, x, y, 6.05, 0.09, RULE_COLS[i])
    # Code pill
    rect(sl, x + 0.15, y + 0.22, 0.62, 0.4, RULE_COLS[i])
    txbox(sl, code, x + 0.17, y + 0.24, 0.58, 0.34,
          bold=True, size=15, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, name, x + 0.88, y + 0.22, 5.0, 0.42,
          bold=True, size=16, color=DARK_TEXT)
    txbox(sl, desc, x + 0.18, y + 0.75, 5.65, 0.75,
          size=14, color=MED_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ML Anomaly Detection
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Task 4: ML Anomaly Detection — Isolation Forest")
two_col(sl,
    "Feature Engineering",
    [
        "Amount (log-scaled to correct skew)",
        "Day of week & month (cyclic encoding)",
        "Voucher type (label-encoded)",
        "Narration length (text quality proxy)",
        "Ledger frequency (account activity score)",
    ],
    "Model & Output",
    [
        "Algorithm: Isolation Forest (unsupervised)",
        "Contamination: configurable (default 5%)",
        "Trained live on each uploaded GL batch",
        "Adds ml_anomaly_flag column to output",
        "Model persisted as models/iforest.pkl",
    ]
)
# Insight callout box
rect(sl, 0.4, 5.95, 12.45, 1.15, RGBColor(0xE8, 0xF8, 0xFF), CYAN, 1.2)
rect(sl, 0.4, 5.95, 0.08, 1.15, CYAN)
txbox(sl, "Key Insight", 0.6, 5.98, 3.0, 0.38, bold=True, size=15, color=CYAN)
txbox(sl, "ML catches statistical outliers that rule-based checks miss — "
          "giving auditors two independent, complementary detection layers.",
      0.6, 6.38, 12.0, 0.65, size=14, italic=True, color=DARK_TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Pipeline & Multi-Agent
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Task 5 & 6: Pipeline & Multi-Agent Architecture")
two_col(sl,
    "Master Pipeline (pipeline.py)",
    [
        "Step 1: Ingest & validate GL schema",
        "Step 2: Run R1-R6 rule engine",
        "Step 3: Train + run Isolation Forest",
        "Step 4: Export flagged Excel report",
        "CLI flags: --input, --output, --no-ml",
    ],
    "Multi-Agent Orchestration",
    [
        "ValidatorAgent -- schema check & cleaning",
        "DetectorAgent -- R1-R6 rules + Isolation Forest",
        "ReporterAgent -- Excel export + summary stats",
        "Orchestrator -- chains agents, collects logs",
        "Each agent has isolated, testable responsibility",
    ]
)
# Flow bar
rect(sl, 0.4, 6.55, 12.45, 0.56, NAVY)
txbox(sl, "GL File  →  ValidatorAgent  →  DetectorAgent  →  ReporterAgent  →  Excel Report",
      0.5, 6.62, 12.0, 0.42, bold=True, size=15, color=CYAN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FastAPI Backend
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Task 6: FastAPI REST Backend")
two_col(sl,
    "API Endpoints",
    [
        "POST /generate — generate synthetic GL",
        "POST /scrutiny — run full scrutiny pipeline",
        "Returns flagged rows + summary statistics",
        "CORS middleware for React frontend integration",
    ],
    "Architecture",
    [
        "Pydantic schemas for request/response validation",
        "Service layer separates business logic from routing",
        "Routers: generator.py + scrutiny.py",
        "Schemas: generator.py + scrutiny.py",
        "Decoupled design — serves React frontend",
    ]
)
rect(sl, 0.4, 5.95, 12.45, 0.48, NAVY_MID)
txbox(sl, "Stack: FastAPI + Uvicorn + Pydantic  |  backend/main.py wires all routers together",
      0.6, 6.02, 12.0, 0.36, size=14, italic=True, color=CYAN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — React / TypeScript Frontend
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Task 7: React / TypeScript Frontend")
two_col(sl,
    "Generator Tab",
    [
        "Config panel: date range, rows, seed",
        "6 anomaly toggle switches",
        "DataTable preview of generated GL",
        "Download CSV button",
    ],
    "Scrutiny Tab",
    [
        "File upload (CSV/XLSX) with drag-and-drop",
        "Live filters: category, date, amount range",
        "Category breakdown bar chart",
        "Summary cards: Total, Flagged, Rule %, ML %",
        "Flagged transactions DataTable + export",
    ]
)
rect(sl, 0.4, 5.95, 12.45, 0.48, NAVY_MID)
txbox(sl, "Stack: React 18 + TypeScript + Vite + Axios  |  Fully typed API layer  |  Component-level state with hooks",
      0.6, 6.02, 12.0, 0.36, size=14, italic=True, color=CYAN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Understanding of Tasks
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "My Understanding of the Tasks")
dot_bullet(sl, [
    "Understood how audit rule-engines map to real-world ICAI GL scrutiny guidelines",
    "Learned why unsupervised ML is preferred for audit — no labelled fraud dataset exists",
    "Understood importance of schema validation at ingestion to prevent downstream failures",
    "Learned how the digital twin approach enables realistic testing without exposing sensitive data",
    "Grasped multi-agent architecture — separates concerns, easier to test and extend independently",
    "Understood REST API design patterns: routers, schemas, services separation in FastAPI",
    "Learned React component architecture: API layer, type definitions, live filters and charts",
], top=1.68, size=17, spacing=0.55)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Impact
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Impact")

IMPACT = [
    ("7 / 7", "Tasks Completed",         "All assigned tasks done end-to-end"),
    ("1,769+", "Lines of Code",           "Across Python & TypeScript"),
    ("6",      "Audit Rules",             "Covering standard GL anomaly patterns (R1-R6)"),
    ("2",      "Detection Layers",        "Rule engine + Isolation Forest ML"),
    ("2",      "Full-Stack Interfaces",   "FastAPI backend + React frontend"),
    ("50,000+","Rows per Run",            "Synthetic GL generated for realistic testing"),
]

for i, (num, label, desc) in enumerate(IMPACT):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.28
    y = 1.7 + row * 2.45
    rect(sl, x, y, 4.0, 2.2, LIGHT_GRAY, CARD_BORDER, 0.8)
    rect(sl, x, y, 4.0, 0.08, CYAN if i % 2 == 0 else ORANGE)
    txbox(sl, num,   x + 0.2, y + 0.22, 3.5, 0.75, bold=True, size=36,
          color=NAVY, align=PP_ALIGN.CENTER)
    txbox(sl, label, x + 0.2, y + 0.95, 3.5, 0.42, bold=True, size=15,
          color=DARK_TEXT, align=PP_ALIGN.CENTER)
    txbox(sl, desc,  x + 0.2, y + 1.38, 3.5, 0.65, size=13,
          color=MED_GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Future Scope
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, WHITE)
header_bar(sl, "Future Scope")
dot_bullet(sl, [
    "Integrate LLM-based narration analysis for deeper weak-narration detection",
    "Add supervised ML model (XGBoost) once labelled audit datasets are available",
    "Extend rules to cover GSTR reconciliation and TDS anomaly patterns",
    "Deploy as SaaS with multi-tenant support for CA firms",
    "Add SHAP explainability layer to justify ML anomaly decisions to auditors",
    "Automate scheduled scrutiny runs with email / Slack alert integration",
], top=1.68, size=18, spacing=0.6)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, NAVY)
rect(sl, 0, 0, 13.33, 7.5, NAVY)
rect(sl, 0, 3.65, 13.33, 0.07, CYAN)
rect(sl, 0, 0, 0.12, 7.5, ORANGE)
rect(sl, 13.21, 0, 0.12, 7.5, CYAN)

txbox(sl, "Thank You", 1.0, 2.2, 11.33, 1.3,
      bold=True, size=58, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Abhinaya  |  RCTS, IIIT-H",
      1.0, 3.85, 11.33, 0.6, size=20, color=CYAN, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
OUT = r"c:\Users\abhin\OneDrive\Desktop\audit-master (1)\audit-master\audit\Abhinaya_Internship_Review_v3.pptx"
prs.save(OUT)
print("Saved:", OUT)
