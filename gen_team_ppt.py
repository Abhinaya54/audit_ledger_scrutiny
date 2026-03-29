"""
gen_team_ppt.py
Generates a detailed team status presentation for the Ledger Scrutiny Assistant project.
Run: python gen_team_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette (matches the teal UI) ────────────────────────────────────
TEAL        = RGBColor(0x0F, 0x76, 0x6E)   # #0F766E  — primary
TEAL_DARK   = RGBColor(0x13, 0x4E, 0x4A)   # #134E4A  — sidebar / dark bg
TEAL_MID    = RGBColor(0x14, 0xB8, 0xA6)   # #14B8A6  — accent / gradient end
TEAL_LIGHT  = RGBColor(0xCC, 0xFB, 0xF1)   # teal-50
TEAL_SOFT   = RGBColor(0x99, 0xF6, 0xE4)   # teal-200
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OFF_WHITE   = RGBColor(0xF8, 0xFA, 0xFC)
LIGHT_GRAY  = RGBColor(0xF1, 0xF5, 0xF9)
MED_GRAY    = RGBColor(0x64, 0x74, 0x8B)
DARK_TEXT   = RGBColor(0x1E, 0x29, 0x3B)
BLACK_SOFT  = RGBColor(0x0F, 0x17, 0x2A)
RED         = RGBColor(0xDC, 0x26, 0x26)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
GREEN_BG    = RGBColor(0xD1, 0xFA, 0xE5)
GREEN_TXT   = RGBColor(0x06, 0x5F, 0x46)
AMBER       = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE      = RGBColor(0x7C, 0x3A, 0xED)
SLATE_200   = RGBColor(0xE2, 0xE8, 0xF0)
CARD_BORDER = RGBColor(0xCB, 0xD5, 0xE1)
ROW_A       = RGBColor(0xF0, 0xFD, 0xFA)   # teal-50
ROW_B       = RGBColor(0xF8, 0xFA, 0xFC)

RULE_COLS   = [TEAL, ORANGE, RED, PURPLE, GREEN, AMBER]

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)


# ── Helpers ──────────────────────────────────────────────────────────────────

def set_bg(slide, color=OFF_WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill_color, line_color=None, line_pt=0.75, radius=0):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color:
        sh.line.color.rgb = line_color
        sh.line.width = Pt(line_pt)
    else:
        sh.line.fill.background()
    return sh


def txbox(slide, text, left, top, width, height,
          bold=False, italic=False, size=14,
          color=None, align=PP_ALIGN.LEFT, wrap=True):
    if color is None:
        color = DARK_TEXT
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold   = bold
    run.font.italic = italic
    run.font.size   = Pt(size)
    run.font.color.rgb = color
    return tb


def header_bar(slide, title, subtitle=None):
    """Teal dark top bar with white title + teal accent strip."""
    rect(slide, 0, 0, 13.33, 1.42, TEAL_DARK)
    rect(slide, 0, 1.42, 13.33, 0.06, TEAL_MID)
    txbox(slide, title, 0.48, 0.14, 12.2, 0.82,
          bold=True, size=28, color=WHITE)
    if subtitle:
        txbox(slide, subtitle, 0.48, 0.92, 12.2, 0.4,
              size=14, color=TEAL_SOFT)


def dot_bullet(slide, items, left=0.55, top=1.7, width=12.0,
               size=16, spacing=0.5, color=None):
    if color is None:
        color = DARK_TEXT
    for i, item in enumerate(items):
        is_sub = item.startswith("    ")
        dot_color = MED_GRAY if is_sub else TEAL
        dot = "–" if is_sub else "●"
        txbox(slide, dot, left, top + i * spacing, 0.3, 0.42,
              size=size if not is_sub else size - 1, color=dot_color, bold=(not is_sub))
        txbox(slide, item.strip(), left + 0.28, top + i * spacing,
              width - 0.28, 0.42, size=size if not is_sub else size - 1, color=color)


def two_col(slide, t1, items1, t2, items2, top=1.62):
    rect(slide, 0.4, top, 5.9, 0.40, TEAL_DARK)
    txbox(slide, t1, 0.55, top + 0.04, 5.6, 0.34,
          bold=True, size=16, color=TEAL_MID)
    for i, item in enumerate(items1):
        txbox(slide, "●  " + item, 0.55, top + 0.52 + i * 0.44, 5.6, 0.42,
              size=14, color=DARK_TEXT)
    rect(slide, 6.95, top, 5.9, 0.40, TEAL_DARK)
    txbox(slide, t2, 7.08, top + 0.04, 5.6, 0.34,
          bold=True, size=16, color=TEAL_MID)
    for i, item in enumerate(items2):
        txbox(slide, "●  " + item, 7.08, top + 0.52 + i * 0.44, 5.6, 0.42,
              size=14, color=DARK_TEXT)


def pill(slide, text, x, y, color, text_color=WHITE, size=12):
    rect(slide, x, y, len(text) * 0.09 + 0.3, 0.32, color)
    txbox(slide, text, x + 0.08, y + 0.03, len(text) * 0.09 + 0.18, 0.28,
          bold=True, size=size, color=text_color)


def stat_card(slide, x, y, number, label, desc, top_color=TEAL):
    rect(slide, x, y, 3.9, 2.1, LIGHT_GRAY, CARD_BORDER, 0.6)
    rect(slide, x, y, 3.9, 0.08, top_color)
    txbox(slide, number, x + 0.15, y + 0.18, 3.6, 0.72,
          bold=True, size=34, color=TEAL_DARK, align=PP_ALIGN.CENTER)
    txbox(slide, label, x + 0.15, y + 0.90, 3.6, 0.38,
          bold=True, size=14, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    txbox(slide, desc,  x + 0.15, y + 1.32, 3.6, 0.60,
          size=11, color=MED_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, TEAL_DARK)

rect(sl, 0, 0, 13.33, 7.5, TEAL_DARK)
rect(sl, 0, 5.55, 13.33, 0.07, TEAL_MID)
rect(sl, 0, 5.62, 13.33, 1.88, BLACK_SOFT)
rect(sl, 0, 0, 0.14, 7.5, TEAL_MID)
rect(sl, 13.19, 0, 0.14, 7.5, ORANGE)

txbox(sl, "Ledger Scrutiny Assistant",
      0.8, 1.35, 11.5, 1.2, bold=True, size=46, color=WHITE)
txbox(sl, "AI-Powered Audit Intelligence Suite — Project Status Presentation",
      0.8, 2.62, 11.5, 0.65, size=20, color=TEAL_SOFT)

txbox(sl, "● Full-Stack Audit Automation   ● 6 Audit Rules + Isolation Forest ML   ● Enterprise React UI",
      0.8, 3.5, 11.5, 0.5, size=14, color=RGBColor(0x99, 0xF6, 0xE4))

txbox(sl, "Prepared by  Abhinaya  |  RCTS, IIIT-H",
      0.8, 5.75, 7.0, 0.5, size=17, color=LIGHT_GRAY)
txbox(sl, "March 2026",
      10.5, 5.75, 2.5, 0.5, size=15, color=TEAL_SOFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Project Overview", "What is the Ledger Scrutiny Assistant?")

dot_bullet(sl, [
    "An end-to-end Computer-Assisted Audit Technique (CAAT) platform for General Ledger analysis",
    "Auditors upload a GL export file → the system detects suspicious transactions automatically",
    "Two independent detection layers provide comprehensive coverage:",
    "    Rule Engine  —  6 coded audit rules (R1–R6) covering standard GL anomaly patterns",
    "    AI Detection  —  Isolation Forest ML flags statistical outliers beyond rule reach",
    "Results surfaced through a professional 5-page React UI with charts, risk scoring & Excel export",
    "Synthetic GL generator enables safe testing without exposing real client data",
], top=1.65, size=17, spacing=0.54)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "System Architecture", "End-to-end data flow from GL upload to audit report")

# Three main blocks
BLOCKS = [
    (0.35,  "FRONTEND\n(React + TS)",    "Upload page\nDashboard\nResults table\nInsights report\nData generator", TEAL),
    (4.75,  "BACKEND\n(FastAPI)",         "POST /scrutiny\nPOST /generate\nPydantic schemas\nCORS middleware", ORANGE),
    (9.15,  "SCRUTINY ENGINE\n(Python)",  "Ingest & validate\nR1–R6 rules\nIsolation Forest\nExcel export", PURPLE),
]

for bx, (bx_x, title_txt, body_txt, col) in enumerate(BLOCKS):
    rect(sl, bx_x, 1.65, 3.65, 4.85, LIGHT_GRAY, CARD_BORDER, 0.8)
    rect(sl, bx_x, 1.65, 3.65, 0.55, col)
    txbox(sl, title_txt, bx_x + 0.15, 1.68, 3.35, 0.50,
          bold=True, size=14, color=WHITE)
    for li, line in enumerate(body_txt.split('\n')):
        txbox(sl, "►  " + line, bx_x + 0.18, 2.35 + li * 0.52, 3.28, 0.46,
              size=13, color=DARK_TEXT)
    # Arrow between blocks
    if bx < 2:
        rect(sl, bx_x + 3.65, 3.75, 1.1, 0.09, TEAL_MID)
        txbox(sl, "▶", bx_x + 4.55, 3.65, 0.4, 0.32,
              size=14, color=TEAL_DARK, bold=True)

# Bottom: data store + output
rect(sl, 0.35, 6.65, 12.6, 0.52, TEAL_DARK)
txbox(sl, "GL File (CSV / XLSX)  ──►  REST API  ──►  Rules + ML Engine  ──►  Flagged Rows + Summary JSON  ──►  React UI + Excel Report",
      0.5, 6.70, 12.0, 0.40, bold=True, size=13, color=TEAL_SOFT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Backend Module Map
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Backend — Module Structure", "Python packages inside audit/backend/")

MODULES = [
    ("main.py",             "FastAPI app wiring, CORS, router registration",                         TEAL),
    ("routers/scrutiny.py", "POST /scrutiny  →  ingest → rules → ML → return JSON",                 TEAL),
    ("routers/generator.py","POST /generate  →  build config → generate GL → return CSV/JSON",       TEAL),
    ("scrutiny/engine.py",  "Orchestrates rule-passes and ML pipeline per uploaded file",             ORANGE),
    ("scrutiny/ingestor.py","Schema validation, dtype normalisation, error reporting",                ORANGE),
    ("scrutiny/rules/",     "R1 Round  ·  R2 Weekend  ·  R3 Period-End  ·  R4 Narration  ·  R5 Dup  ·  R6 Manual", ORANGE),
    ("scrutiny/ml/",        "feature_engineering.py  +  model.py (IsolationForest train & predict)",  PURPLE),
    ("scrutiny/exporter.py","Builds Excel report with flagged rows + summary sheet",                  PURPLE),
    ("generator/core.py",   "Generates 50 000+ GL rows with configurable params & anomaly injection", GREEN),
    ("agents/",             "ValidatorAgent · DetectorAgent · ReporterAgent · Orchestrator",          RED),
]

for i, (mod, desc, col) in enumerate(MODULES):
    col_n = i % 2
    row_n = i // 2
    bx = 0.35 + col_n * 6.5
    by = 1.62 + row_n * 0.54
    rect(sl, bx, by, 6.1, 0.46, LIGHT_GRAY, CARD_BORDER, 0.5)
    rect(sl, bx, by, 0.06, 0.46, col)
    txbox(sl, mod, bx + 0.14, by + 0.03, 2.1, 0.38,
          bold=True, size=12, color=col)
    txbox(sl, desc, bx + 2.28, by + 0.04, 3.7, 0.38,
          size=11, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Audit Rules Engine
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Audit Rules Engine — R1 to R6", "6 rule modules in scrutiny/rules/ — each independently testable")

RULES = [
    ("R1", "Round Numbers",    "Flags amounts that are multiples of 1,000 or 10,000 — indicates estimated or manual entries",            TEAL),
    ("R2", "Weekend Entries",  "Transactions posted on Sundays — may indicate unauthorised out-of-cycle postings",                      ORANGE),
    ("R3", "Period-End",       "Entries within the critical 5-day period-end window — classic revenue management red flag",              RED),
    ("R4", "Weak Narration",   "Description too short (< 5 chars) or generic placeholder — breaks audit trail compliance",              PURPLE),
    ("R5", "Duplicate Check",  "Same date + ledger + amount appearing multiple times — potential duplicate payment",                    GREEN),
    ("R6", "Manual Journal",   "Manual JV voucher type with no supporting narration — requires additional source verification",          AMBER),
]

for i, (code, name, desc, col) in enumerate(RULES):
    col_pos = i % 2
    row_pos = i // 2
    x = 0.35 + col_pos * 6.5
    y = 1.62 + row_pos * 1.88
    rect(sl, x, y, 6.1, 1.72, LIGHT_GRAY, CARD_BORDER, 0.7)
    rect(sl, x, y, 6.1, 0.08, col)
    rect(sl, x + 0.15, y + 0.22, 0.65, 0.38, col)
    txbox(sl, code, x + 0.17, y + 0.24, 0.60, 0.32,
          bold=True, size=14, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, name, x + 0.92, y + 0.20, 5.0, 0.40,
          bold=True, size=15, color=DARK_TEXT)
    txbox(sl, desc, x + 0.18, y + 0.72, 5.72, 0.82,
          size=12, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ML Anomaly Detection
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "AI Anomaly Detection — Isolation Forest", "Unsupervised ML layer in scrutiny/ml/")

two_col(sl,
    "Feature Engineering  (feature_engineering.py)",
    [
        "amount_log  —  log-scaled to correct right skew",
        "day_of_week  —  sin/cos cyclic encoding",
        "month  —  sin/cos cyclic encoding",
        "voucher_type_enc  —  label-encoded",
        "narration_len  —  text quality proxy",
        "ledger_freq  —  account activity score",
    ],
    "Model & Output  (model.py)",
    [
        "Algorithm: sklearn IsolationForest (unsupervised)",
        "Contamination: configurable per request (0.02–0.18)",
        "Trained fresh on each uploaded GL batch",
        "Adds ml_anomaly_flag (–1 = anomaly) to output",
        "Adds ml_anomaly_score for risk prioritisation",
        "Model saved as backend/models/iforest.pkl",
    ],
    top=1.62,
)

rect(sl, 0.35, 5.88, 12.6, 1.28, TEAL_LIGHT, TEAL_SOFT, 1.0)
rect(sl, 0.35, 5.88, 0.08, 1.28, TEAL)
txbox(sl, "Why Isolation Forest?", 0.55, 5.92, 4.0, 0.36,
      bold=True, size=14, color=TEAL)
txbox(sl,
      "No labelled fraud dataset exists for audit GL data — supervised models are not feasible. "
      "Isolation Forest is ideal because it learns \"what is normal\" from the data itself and flags "
      "statistical outliers that rule-based checks will never catch. "
      "The two layers are complementary: rules handle known patterns, ML handles unknown ones.",
      0.55, 6.30, 12.2, 0.78, size=13, italic=True, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Synthetic GL Generator
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Synthetic GL Generator", "generator/ module — privacy-safe test data creation")

two_col(sl,
    "Core Generator",
    [
        "Generates realistic GL entries — 50,000+ rows per run",
        "Configurable: fiscal year, accounts, voucher types",
        "Randomised narrations using Faker (en_IN locale)",
        "CLI flags: --period, --rows, --seed, --no-inject",
        "Injects all 6 anomaly types at configurable rates",
    ],
    "Digital Twin Mode",
    [
        "Mirrors real GL statistics without raw data exposure",
        "Input: JSON summary (mean, std, distributions)",
        "Preserves amount, weekday & account distributions",
        "Privacy-safe  —  no original transactions in output",
        "Ideal for CA firms to share profiles without data",
    ],
    top=1.62,
)

# Anomaly injection bar
rect(sl, 0.35, 5.55, 12.6, 0.44, TEAL_DARK)
txbox(sl, "Anomaly Types Injected:", 0.5, 5.59, 2.5, 0.36,
      bold=True, size=13, color=TEAL_MID)
txbox(sl, "Round Numbers  |  Weekend Entries  |  Period-End Spikes  |  Weak Narrations  |  Duplicates  |  Manual Journals",
      3.1, 5.59, 9.8, 0.36, size=13, color=LIGHT_GRAY)

rect(sl, 0.35, 6.08, 12.6, 1.15, LIGHT_GRAY, CARD_BORDER, 0.6)
rect(sl, 0.35, 6.08, 0.06, 1.15, GREEN)
txbox(sl, "Frontend Integration", 0.55, 6.12, 3.5, 0.32, bold=True, size=13, color=GREEN)
txbox(sl,
      "The Data Generator tab in the React UI calls POST /generate, lets auditors configure "
      "all params via toggles and sliders, previews the generated GL in a DataTable, "
      "and provides a one-click CSV download — no CLI required.",
      0.55, 6.46, 12.1, 0.68, size=12, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FastAPI Backend
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "FastAPI REST Backend", "backend/main.py + routers/scrutiny.py + routers/generator.py")

two_col(sl,
    "API Endpoints",
    [
        "POST /api/scrutiny  —  run full audit on uploaded GL file",
        "POST /api/generate  —  generate synthetic GL with params",
        "GET  /api/scrutiny/export  —  download flagged Excel report",
        "Returns: flagged_rows[], category_counts[], summary{} JSON",
        "Pydantic models enforce request/response type safety",
    ],
    "Architecture & Dev Setup",
    [
        "Run from audit/backend/:  uvicorn main:app --reload",
        "CORS enabled for http://localhost:5173 (Vite dev server)",
        "Vite proxy: /api  →  http://localhost:8000 (no CORS in prod)",
        "Multipart file upload via python-multipart",
        "Service layer separates business logic from routing",
    ],
    top=1.62,
)

# Response schema box
rect(sl, 0.35, 5.55, 12.6, 1.65, RGBColor(0x1E, 0x29, 0x3B), None)
txbox(sl, "Response Schema (ScrutinyResponse)", 0.5, 5.60, 5.0, 0.35,
      bold=True, size=13, color=TEAL_MID)
txbox(sl,
      "{ flagged_rows: [ { date, ledger_name, amount, narration, scrutiny_category, "
      "ml_anomaly_flag, ml_anomaly_score } ],\n"
      "  category_counts: [ { category, count } ],\n"
      "  summary: { total_entries, total_flagged, rule_flagged, ml_flagged, pct_flagged } }",
      0.5, 5.98, 12.1, 1.12, size=11, color=TEAL_SOFT, italic=True)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Frontend UI Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Enterprise React Frontend — 5 Pages", "React 19 + TypeScript + Vite + Tailwind CSS v4 + Recharts")

PAGES = [
    ("Dashboard",            "KPI cards, anomaly trend line chart, risk pie, breakdown bars, recent suspicious feed, CTA banner",  TEAL,   "●"),
    ("Upload & Analyze",     "File drag-and-drop (CSV/XLSX), AI toggle, 5-step risk-coloured sensitivity slider, Run Analysis CTA", ORANGE, "▲"),
    ("Suspicious Transactions","Custom risk table (200 rows), risk score badges, sort/filter, bar + pie charts, Excel export",      RED,    "⚠"),
    ("Insights & Report",    "AI-generated insights, anomaly trend, breakdown bars, top accounts, full Excel download",            PURPLE, "📄"),
    ("Data Generator",       "Config panel, anomaly toggles, DataTable preview of synthetic GL, download CSV",                     GREEN,  "⚙"),
]

for i, (name, features, col, icon) in enumerate(PAGES):
    y = 1.62 + i * 1.02
    rect(sl, 0.35, y, 12.6, 0.90, LIGHT_GRAY, CARD_BORDER, 0.6)
    rect(sl, 0.35, y, 0.06, 0.90, col)
    # Number badge
    rect(sl, 0.48, y + 0.18, 0.52, 0.52, col)
    txbox(sl, str(i + 1), 0.48, y + 0.18, 0.52, 0.52,
          bold=True, size=18, color=WHITE, align=PP_ALIGN.CENTER)
    txbox(sl, name, 1.12, y + 0.06, 3.2, 0.38,
          bold=True, size=14, color=col)
    txbox(sl, features, 1.12, y + 0.46, 11.6, 0.36,
          size=12, color=MED_GRAY)

rect(sl, 0.35, 6.77, 12.6, 0.46, TEAL_DARK)
txbox(sl, "State management lifted to AppShell.tsx  ·  Vite proxy to FastAPI  ·  TypeScript strict mode — zero any types  ·  Recharts for all visuals",
      0.5, 6.82, 12.0, 0.36, size=12, italic=True, color=TEAL_SOFT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Dashboard Page
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Page 1: Dashboard", "Real-time audit overview — auto-populates after every analysis run")

two_col(sl,
    "KPI Cards (4 metrics)",
    [
        "Total Transactions — full ledger entry count",
        "Suspicious Transactions — count + % of total",
        "AI Detections — ML-only flagged count",
        "Risk Level — Low / Medium / High (dynamic colour)",
    ],
    "Charts & Activity",
    [
        "Anomaly Trend — monthly line chart over time",
        "Anomaly Distribution — clean vs suspicious pie",
        "Breakdown — category bar chart (top 5 rules)",
        "Recent Suspicious Activity — last 5 flagged rows",
    ],
    top=1.62,
)

rect(sl, 0.35, 5.55, 12.6, 1.62, LIGHT_GRAY, CARD_BORDER, 0.6)
rect(sl, 0.35, 5.55, 0.06, 1.62, TEAL)
txbox(sl, "CTA Banner", 0.55, 5.59, 2.8, 0.34, bold=True, size=13, color=TEAL)
txbox(sl,
      "A teal banner at the bottom shows the total flagged count and surfaces two action buttons: "
      "\"Re-analyze\" (returns to Upload page) and \"Review Results\" (jumps to Suspicious Transactions). "
      "The empty state guides new users with a feature grid explaining Rule Checks, AI Detection, and Risk Scoring.",
      0.55, 5.95, 12.1, 1.12, size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Upload & Analyze Page
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Page 2: Upload & Analyze", "Step-based workflow — file validation + configurable analysis settings")

two_col(sl,
    "Step 1 — File Upload",
    [
        "Drag-and-drop zone with visual feedback states",
        "Validates extension (CSV/XLSX) + MIME type",
        "Enforces 50 MB size limit with clear error messages",
        "Shows file name + size on successful selection",
        "Click-to-replace replaces without page reload",
    ],
    "Step 2 — Analysis Settings",
    [
        "AI Toggle — enables/disables Isolation Forest pass",
        "Detection Sensitivity — 5 risk-coded steps:",
        "    Very Low (gray)  ·  Low (green)  ·  Balanced (teal)",
        "    High (orange)  ·  Very High (red)",
        "Maps visually: low sensitivity = low risk colour",
    ],
    top=1.62,
)

rect(sl, 0.35, 5.55, 12.6, 1.62, LIGHT_GRAY, CARD_BORDER, 0.6)
rect(sl, 0.35, 5.55, 0.06, 1.62, ORANGE)
txbox(sl, "Sensitivity → Contamination mapping", 0.55, 5.59, 4.5, 0.34, bold=True, size=13, color=ORANGE)
SENS = [
    ("Very Low", "0.02", RGBColor(0x64, 0x74, 0x8B)),
    ("Low",      "0.04", GREEN),
    ("Balanced", "0.07", TEAL),
    ("High",     "0.12", AMBER),
    ("Very High","0.18", RED),
]
for i, (label, val, col) in enumerate(SENS):
    bx = 0.55 + i * 2.42
    rect(sl, bx, 6.00, 2.2, 0.75, LIGHT_GRAY, CARD_BORDER, 0.5)
    rect(sl, bx, 6.00, 2.2, 0.06, col)
    txbox(sl, label, bx + 0.1, 6.10, 2.0, 0.28, bold=True, size=12, color=col)
    txbox(sl, "contamination = " + val, bx + 0.1, 6.38, 2.0, 0.28, size=10, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Flagged Transactions Page
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Page 3: Suspicious Transactions", "Risk-scored, sortable, filterable table of all flagged entries")

two_col(sl,
    "Table Features",
    [
        "7 columns: Date · Account · Amount · Narration · Type · Risk · Flag",
        "Risk Score (0–100): computed from severity + ML score",
        "Row colour: High=red tint, Medium=amber tint, Low=white",
        "Sort by Risk Score toggle — top risks surface instantly",
        "Show High Risk Only toggle — ≥75 score filter",
        "Renders up to 200 rows inline; full data via Excel export",
    ],
    "Charts & Filters",
    [
        "Category bar chart — top 6 anomaly types by count",
        "Severity pie chart — High / Medium / Low distribution",
        "LiveFilters: rule type, severity, account, date range",
        "Active filter count badge updates in real-time",
        "\"View Insights & Report\" CTA transitions to Page 4",
        "Export Report button triggers direct Excel download",
    ],
    top=1.62,
)

rect(sl, 0.35, 5.55, 12.6, 1.62, LIGHT_GRAY, CARD_BORDER, 0.6)
rect(sl, 0.35, 5.55, 0.06, 1.62, PURPLE)
txbox(sl, "Risk Score Computation", 0.55, 5.59, 3.5, 0.34, bold=True, size=13, color=PURPLE)
txbox(sl,
      "Base score from severity:  High → 82  ·  Medium → 57  ·  Low → 32\n"
      "If ML flagged (flag = –1):  abs(ml_anomaly_score) × 200, capped at 100, takes maximum of base vs ML score\n"
      "Badge colour:  score ≥ 75 → red  ·  score ≥ 50 → amber  ·  else → green",
      0.55, 5.96, 12.1, 1.12, size=12, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Insights & Report Page
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Page 4: Insights & Report", "AI-generated audit findings + downloadable Excel report")

two_col(sl,
    "Report Sections",
    [
        "Report Header — date, engine version, risk badge",
        "Executive Summary — 4 KPI cards (entries, rule flags, ML detections, total)",
        "Key Findings — peak month, top category, total suspicious value",
        "AI-Generated Insights — 6–8 data-driven audit observations",
        "Anomaly Trend — monthly line chart",
        "Findings by Type — horizontal bar chart with % share",
        "Top Suspicious Accounts — ranked by flag count",
        "Audit Conclusion — CAAT summary + legal disclaimer",
    ],
    "AI Insights Engine",
    [
        "Computes insights live from real data — no mock text",
        "Covers: anomaly rate assessment (low/med/high)",
        "Period-end concentration warning",
        "Duplicate payment risk alert",
        "Round-number entry observation",
        "ML unique detection count",
        "Weekend posting irregularity flag",
        "Total suspicious value with priority guidance",
    ],
    top=1.62,
)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Data Generator Tab
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Page 5: Data Generator Tab", "In-browser synthetic GL generation — no CLI needed")

two_col(sl,
    "Configuration Panel",
    [
        "Date range selector (fiscal year start/end)",
        "Number of rows slider",
        "Random seed input for reproducibility",
        "Company name and base currency",
        "Digital Twin toggle with JSON profile upload",
    ],
    "Anomaly Toggles",
    [
        "6 individual on/off switches per anomaly type",
        "Round Numbers  —  Weekend Entries",
        "Period-End  —  Weak Narrations",
        "Duplicate Entries  —  Manual Journals",
        "Injection rate configurable per type",
    ],
    top=1.62,
)

rect(sl, 0.35, 5.55, 12.6, 1.62, LIGHT_GRAY, CARD_BORDER, 0.6)
rect(sl, 0.35, 5.55, 0.06, 1.62, GREEN)
txbox(sl, "Output Preview", 0.55, 5.59, 2.5, 0.34, bold=True, size=13, color=GREEN)
txbox(sl,
      "Generated GL is shown immediately in a paginated DataTable (50 rows/page). "
      "Columns include: date, ledger_name, amount, voucher_type, narration, debit, credit.\n"
      "A \"Download CSV\" button streams the full dataset to the browser. "
      "The generated file can immediately be uploaded to the Upload & Analyze page for end-to-end testing.",
      0.55, 5.95, 12.1, 1.12, size=13, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Design System & Teal Palette
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "UI Design System — Teal Palette", "Professional fintech aesthetic — calm authority, risk-aware colour coding")

# Palette swatches row
SWATCHES = [
    ("#134E4A", "Sidebar Dark",     TEAL_DARK),
    ("#0F766E", "Primary",          TEAL),
    ("#115E59", "Hover",            RGBColor(0x11, 0x5E, 0x59)),
    ("#14B8A6", "Accent",           TEAL_MID),
    ("#CCFBF1", "Surface",          TEAL_LIGHT),
]
for i, (hex_val, label, col) in enumerate(SWATCHES):
    bx = 0.38 + i * 2.52
    rect(sl, bx, 1.66, 2.3, 1.12, col, CARD_BORDER, 0.5)
    txbox(sl, hex_val, bx + 0.08, 2.82, 2.14, 0.30,
          bold=True, size=11, color=MED_GRAY)
    txbox(sl, label,   bx + 0.08, 3.14, 2.14, 0.28,
          size=12, color=DARK_TEXT, bold=True)

# Risk sensitivity colours
txbox(sl, "Detection Sensitivity Pills — Risk-Matched Colours", 0.38, 3.60, 9.0, 0.36,
      bold=True, size=14, color=DARK_TEXT)
RISK_PILLS = [
    ("Very Low",  RGBColor(0x64, 0x74, 0x8B), WHITE, "Conservative — few flags"),
    ("Low",       GREEN,                       WHITE, "Slight sensitivity"),
    ("Balanced",  TEAL,                        WHITE, "Default setting"),
    ("High",      AMBER,                       WHITE, "More aggressive"),
    ("Very High", RED,                         WHITE, "Maximum recall"),
]
for i, (label, bg, fg, tip) in enumerate(RISK_PILLS):
    bx = 0.38 + i * 2.52
    rect(sl, bx, 4.05, 2.3, 0.44, bg)
    txbox(sl, label, bx + 0.1, 4.08, 2.1, 0.36,
          bold=True, size=13, color=fg, align=PP_ALIGN.CENTER)
    txbox(sl, tip, bx + 0.06, 4.55, 2.2, 0.30,
          size=10, color=MED_GRAY, align=PP_ALIGN.CENTER)

# Component patterns
rect(sl, 0.38, 5.05, 12.6, 1.0, LIGHT_GRAY, CARD_BORDER, 0.5)
rect(sl, 0.38, 5.05, 0.06, 1.0, TEAL)
txbox(sl, "Consistent Component Patterns", 0.58, 5.09, 4.0, 0.32, bold=True, size=13, color=TEAL)
txbox(sl,
      "MetricCard (5 variants: teal/green/red/amber/purple)  ·  SectionCard with numbered steps  ·  "
      "DataTable with pagination + hiddenColumns  ·  LiveFilters with multi-select  ·  "
      "ToggleSwitch  ·  Tooltip  ·  WorkflowSteps indicator  ·  AnalysisOverlay with progress checklist",
      0.58, 5.44, 12.1, 0.52, size=12, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Completion Status Table
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Completion Status — What's Done", "Full feature inventory across backend and frontend")

TASKS = [
    ("Backend",  "Synthetic GL Generator (core + digital twin + CLI)",    "100%", "Done"),
    ("Backend",  "Rule-Based Scrutiny Engine — R1 to R6",                 "100%", "Done"),
    ("Backend",  "ML Anomaly Detection — Isolation Forest",               "100%", "Done"),
    ("Backend",  "Master Pipeline (ingest → rules → ML → export)",        "100%", "Done"),
    ("Backend",  "Multi-Agent System (Validator, Detector, Reporter)",     "100%", "Done"),
    ("Backend",  "FastAPI REST Backend with Pydantic schemas + CORS",      "100%", "Done"),
    ("Frontend", "AppShell + Sidebar + WorkflowSteps + AnalysisOverlay",  "100%", "Done"),
    ("Frontend", "Dashboard — KPIs, trend chart, pie, breakdown, CTA",    "100%", "Done"),
    ("Frontend", "Upload & Analyze — file validation + sensitivity slider","100%", "Done"),
    ("Frontend", "Suspicious Transactions — risk table + charts + export", "100%", "Done"),
    ("Frontend", "Insights & Report — AI insights + report sections",      "100%", "Done"),
    ("Frontend", "Data Generator Tab — config panel + DataTable + CSV",    "100%", "Done"),
    ("Design",   "Enterprise Teal Palette (#0F766E) + risk-coded colours", "100%", "Done"),
]

COL_X = [0.38, 1.30, 9.80, 11.10]
COL_W = [0.84, 8.40, 1.18, 1.88]
HDRS  = ["Area", "Feature", "%", "Status"]

rect(sl, 0.38, 1.62, 12.6, 0.40, TEAL_DARK)
for hdr, cx, cw in zip(HDRS, COL_X, COL_W):
    txbox(sl, hdr, cx + 0.05, 1.65, cw - 0.08, 0.34,
          bold=True, size=12, color=WHITE)

for i, (area, feat, pct, status) in enumerate(TASKS):
    y = 2.05 + i * 0.38
    bg = ROW_A if i % 2 == 0 else ROW_B
    vals = [area, feat, pct, status]
    for j, (txt, cx, cw) in enumerate(zip(vals, COL_X, COL_W)):
        fill = GREEN_BG if j == 3 else bg
        rect(sl, cx, y, cw, 0.36, fill)
        c = GREEN_TXT if j == 3 else (TEAL if j == 0 else DARK_TEXT)
        txbox(sl, txt, cx + 0.05, y + 0.03, cw - 0.08, 0.30,
              size=10, color=c, bold=(j == 3))

rect(sl, 0.38, 7.07, 4.0, 0.28, TEAL_DARK)
txbox(sl, "13 / 13 Features Completed  ✓", 0.48, 7.10, 3.82, 0.22,
      bold=True, size=11, color=TEAL_SOFT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Technical Stack
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Technical Stack", "Languages, frameworks, and libraries used")

STACK = [
    ("BACKEND",  [
        ("Language",     "Python 3.11+"),
        ("Web Framework","FastAPI + Uvicorn"),
        ("Validation",   "Pydantic v2"),
        ("ML",           "scikit-learn  IsolationForest"),
        ("Data",         "pandas + numpy"),
        ("Export",       "openpyxl (Excel .xlsx)"),
        ("Test Data",    "Faker (en_IN locale)"),
        ("Routing",      "APIRouter — scrutiny + generator"),
    ], TEAL),
    ("FRONTEND", [
        ("Language",     "TypeScript 5 (strict mode)"),
        ("Framework",    "React 19 + Vite 6"),
        ("Styling",      "Tailwind CSS v4"),
        ("Charts",       "Recharts  (Line / Bar / Pie)"),
        ("HTTP",         "Axios  via scrutinyApi.ts"),
        ("Icons",        "Inline SVG (zero icon library)"),
        ("Types",        "types/scrutiny.ts  —  FlaggedRow, DisplayRow"),
        ("State",        "useState / useMemo in AppShell"),
    ], ORANGE),
    ("TOOLING",  [
        ("Build",        "Vite  (fast HMR + production bundle)"),
        ("Version Ctrl", "Git  (branch: main)"),
        ("Formatting",   "Tailwind class ordering"),
        ("Dev Server",   "Uvicorn --reload + Vite dev proxy"),
        ("API Docs",     "FastAPI auto Swagger UI  /docs"),
    ], PURPLE),
]

for bi, (area, items, col) in enumerate(STACK):
    bx = 0.35 + bi * 4.33
    bw = 4.15
    rect(sl, bx, 1.62, bw, 0.40, col)
    txbox(sl, area, bx + 0.12, 1.65, bw - 0.2, 0.32,
          bold=True, size=14, color=WHITE)
    for ri, (key, val) in enumerate(items):
        y = 2.08 + ri * 0.54
        rect(sl, bx, y, bw, 0.48, LIGHT_GRAY if ri % 2 == 0 else OFF_WHITE, CARD_BORDER, 0.4)
        txbox(sl, key + ":", bx + 0.1, y + 0.05, 1.4, 0.36,
              bold=True, size=11, color=col)
        txbox(sl, val, bx + 1.5, y + 0.05, bw - 1.6, 0.36,
              size=11, color=DARK_TEXT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Future Scope
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
header_bar(sl, "Future Scope", "Planned enhancements and extensions")

FUTURE = [
    (TEAL,   "LLM Narration Analysis",  "Integrate an LLM to semantically analyse narration text — catching vague descriptions that pass the character-length check but still lack audit trail clarity."),
    (ORANGE, "Supervised ML Model",     "Once labelled audit datasets are collected, add an XGBoost classifier for higher precision — complementing the unsupervised Isolation Forest baseline."),
    (RED,    "GSTR & TDS Rules",        "Extend the rule engine (R7+) to cover GST reconciliation mismatches and TDS deduction anomalies for Indian statutory compliance."),
    (PURPLE, "SHAP Explainability",     "Add SHAP value visualisations so auditors understand exactly which features drove each ML anomaly decision — building trust in the AI layer."),
    (GREEN,  "Scheduled Runs & Alerts", "Automated nightly scrutiny runs with configurable email / Slack alerts when the flagged rate exceeds a threshold."),
    (AMBER,  "Multi-Tenant SaaS",       "Deploy as a cloud SaaS platform with per-CA-firm isolation, role-based access (auditor / manager / admin), and audit log retention."),
]

for i, (col, title, desc) in enumerate(FUTURE):
    row = i // 2
    c   = i % 2
    bx  = 0.35 + c * 6.5
    by  = 1.62 + row * 1.85
    rect(sl, bx, by, 6.1, 1.68, LIGHT_GRAY, CARD_BORDER, 0.6)
    rect(sl, bx, by, 6.1, 0.07, col)
    txbox(sl, title, bx + 0.18, by + 0.20, 5.70, 0.36,
          bold=True, size=14, color=col)
    txbox(sl, desc,  bx + 0.18, by + 0.62, 5.70, 0.92,
          size=12, color=MED_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Thank You / Q&A
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, TEAL_DARK)

rect(sl, 0, 0, 13.33, 7.5, TEAL_DARK)
rect(sl, 0, 5.35, 13.33, 0.07, TEAL_MID)
rect(sl, 0, 5.42, 13.33, 2.08, BLACK_SOFT)
rect(sl, 0, 0, 0.14, 7.5, TEAL_MID)
rect(sl, 13.19, 0, 0.14, 7.5, ORANGE)

txbox(sl, "Thank You", 1.0, 1.55, 11.33, 1.5,
      bold=True, size=60, color=WHITE, align=PP_ALIGN.CENTER)
txbox(sl, "Questions & Discussion", 1.0, 3.15, 11.33, 0.60,
      size=24, color=TEAL_SOFT, align=PP_ALIGN.CENTER)

txbox(sl, "●  Backend: Python · FastAPI · Isolation Forest · pandas · openpyxl",
      1.5, 5.60, 10.33, 0.38, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txbox(sl, "●  Frontend: React 19 · TypeScript · Vite · Tailwind CSS v4 · Recharts",
      1.5, 5.98, 10.33, 0.38, size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
txbox(sl, "Abhinaya  |  RCTS, IIIT-H  |  March 2026",
      1.5, 6.58, 10.33, 0.45, size=16, color=TEAL_SOFT, align=PP_ALIGN.CENTER)


# ── Save ─────────────────────────────────────────────────────────────────────
OUT = r"c:\Users\abhin\OneDrive\Desktop\audit-master (1)\audit-master\audit\Ledger_Scrutiny_Team_Presentation.pptx"
prs.save(OUT)
print("Saved:", OUT)
